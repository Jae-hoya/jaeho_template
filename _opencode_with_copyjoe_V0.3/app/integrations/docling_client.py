import re
from pathlib import Path
from tempfile import TemporaryDirectory
from time import perf_counter
from typing import Any


class DoclingClient:
    def __init__(
        self,
        pdf_ocr_strategy: str = "rapid",
        pdf_ocr_min_chars: int = 180,
        pdf_layout_model_strategy: str = "off",
        pdf_vlm_preset: str = "smoldocling",
        pdf_vlm_device: str = "auto",
        image_processing_strategy: str = "rapid",
        image_vlm_preset: str = "smoldocling",
        image_vlm_device: str = "auto",
    ) -> None:
        self._converter = None
        self._image_rapidocr_converter = None
        self._image_vlm_converter = None
        self._pdf_vlm_converter = None
        self._pdf_rapidocr_converter = None
        self._pdf_easyocr_converter = None
        self._pdf_ocr_engine = "none"
        self._converter_initialized = False
        self._image_rapidocr_initialized = False
        self._image_vlm_initialized = False
        self._pdf_vlm_initialized = False
        self._pdf_rapidocr_initialized = False
        self._pdf_easyocr_initialized = False
        self._pdf_ocr_strategy = self._normalize_pdf_ocr_strategy(pdf_ocr_strategy)
        self._pdf_ocr_min_chars = max(20, int(pdf_ocr_min_chars))
        self._pdf_layout_model_strategy = self._normalize_pdf_layout_model_strategy(pdf_layout_model_strategy)
        self._pdf_vlm_preset = (pdf_vlm_preset or "smoldocling").strip().lower()
        self._pdf_vlm_device = self._normalize_vlm_device(pdf_vlm_device)
        self._image_processing_strategy = self._normalize_image_processing_strategy(image_processing_strategy)
        self._image_vlm_preset = (image_vlm_preset or "smoldocling").strip().lower()
        self._image_vlm_device = self._normalize_vlm_device(image_vlm_device)

    def warm_up(self) -> dict[str, object]:
        started = perf_counter()
        warmed_components: list[str] = []
        failed_components: list[str] = []

        self._ensure_converter()
        if self._converter is not None:
            warmed_components.append("docling_default")
        else:
            failed_components.append("docling_default")

        with TemporaryDirectory(prefix="copyjoe-ocr-warmup-") as temp_dir:
            temp_path = Path(temp_dir)
            image_path = temp_path / "warmup.png"
            pdf_path = temp_path / "warmup.pdf"

            image_ready = self._create_warmup_image(image_path)
            pdf_ready = self._create_warmup_pdf(pdf_path)

            if image_ready:
                self._warm_up_image_components(image_path, warmed_components, failed_components)
            else:
                failed_components.append("warmup_image_file")

            if pdf_ready:
                self._warm_up_pdf_components(pdf_path, warmed_components, failed_components)
            else:
                failed_components.append("warmup_pdf_file")

        duration_ms = int((perf_counter() - started) * 1000)
        return {
            "ok": len(failed_components) == 0,
            "image_processing_strategy": self._image_processing_strategy,
            "image_vlm_preset": self._image_vlm_preset,
            "image_vlm_device": self._image_vlm_device,
            "pdf_ocr_strategy": self._pdf_ocr_strategy,
            "pdf_layout_model_strategy": self._pdf_layout_model_strategy,
            "pdf_vlm_preset": self._pdf_vlm_preset,
            "pdf_vlm_device": self._pdf_vlm_device,
            "warmed_components": warmed_components,
            "failed_components": failed_components,
            "duration_ms": duration_ms,
        }

    def convert_to_text(self, file_path: Path) -> str:
        text, _ = self.convert_to_text_with_meta(file_path)
        return text

    def convert_to_text_with_meta(self, file_path: Path) -> tuple[str, dict[str, object]]:
        suffix = file_path.suffix.lower()

        if suffix == ".pdf":
            pypdf_text, snapshot = self._extract_pdf_snapshot(file_path)
            if self._should_use_pdf_vlm(snapshot):
                self._ensure_pdf_vlm_converter()
                if self._pdf_vlm_converter is not None:
                    vlm_text = self._convert_with_docling(self._pdf_vlm_converter, file_path).strip()
                    if vlm_text:
                        return vlm_text, {"conversion_engine": "pdf_smoldocling"}
            selected_pdf_text = self._convert_pdf_with_ocr_policy(pypdf_text, file_path)
            if selected_pdf_text:
                return selected_pdf_text, {"conversion_engine": f"pdf_{self._pdf_ocr_engine}"}
            return self._fallback_convert(file_path), {"conversion_engine": "fallback_pdf"}

        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            image_text, image_engine = self._convert_image_with_policy(file_path)
            return image_text, {"conversion_engine": image_engine}

        self._ensure_converter()
        docling_text = self._convert_with_docling(self._converter, file_path).strip()

        if docling_text:
            return docling_text, {"conversion_engine": "docling_default"}

        return self._fallback_convert(file_path), {"conversion_engine": self._fallback_engine_for_suffix(suffix)}

    def _normalize_pdf_ocr_strategy(self, strategy: str) -> str:
        normalized = (strategy or "rapid").strip().lower()
        if normalized in {"off", "rapid", "easy", "hybrid"}:
            return normalized
        return "rapid"

    def _normalize_image_processing_strategy(self, strategy: str) -> str:
        normalized = (strategy or "rapid").strip().lower()
        if normalized in {"rapid", "smoldocling", "hybrid"}:
            return normalized
        return "rapid"

    def _normalize_pdf_layout_model_strategy(self, strategy: str) -> str:
        normalized = (strategy or "auto").strip().lower()
        if normalized in {"off", "auto", "smoldocling"}:
            return normalized
        return "auto"

    def _normalize_vlm_device(self, device: str) -> str:
        normalized = (device or "auto").strip().lower()
        if normalized in {"auto", "cpu", "cuda", "mps", "xpu"}:
            return normalized
        return "auto"

    def _should_use_pdf_vlm(self, snapshot: dict[str, int]) -> bool:
        strategy = self._pdf_layout_model_strategy
        if strategy == "off":
            return False
        if strategy == "smoldocling":
            return True

        image_count = int(snapshot.get("image_count", 0))
        table_like_lines = int(snapshot.get("table_like_lines", 0))
        total_lines = max(1, int(snapshot.get("total_lines", 0)))

        if image_count > 0:
            return True
        if table_like_lines >= 6 and (table_like_lines / total_lines) >= 0.2:
            return True
        return False

    def _extract_pdf_snapshot(self, file_path: Path) -> tuple[str, dict[str, int]]:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            text_parts: list[str] = []
            image_count = 0
            table_like_lines = 0
            total_lines = 0

            for page in reader.pages:
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    text_parts.append(page_text)
                    lines = [line for line in page_text.splitlines() if line.strip()]
                    total_lines += len(lines)
                    table_like_lines += self._estimate_table_like_lines(lines)

                image_count += self._count_pdf_page_images(page)

            return "\n".join(text_parts).strip(), {
                "pages": len(reader.pages),
                "image_count": image_count,
                "table_like_lines": table_like_lines,
                "total_lines": total_lines,
            }
        except Exception:
            return "", {
                "pages": 0,
                "image_count": 0,
                "table_like_lines": 0,
                "total_lines": 0,
            }

    def _estimate_table_like_lines(self, lines: list[str]) -> int:
        score = 0
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            if "|" in stripped and len([part for part in stripped.split("|") if part.strip()]) >= 3:
                score += 1
                continue

            token_count = len(re.findall(r"\S+", stripped))
            numeric_count = len(re.findall(r"\d+(?:[.,]\d+)?", stripped))
            wide_gap_count = len(re.findall(r"\s{2,}", stripped))
            if token_count >= 4 and numeric_count >= 2 and wide_gap_count >= 1:
                score += 1

        return score

    def _count_pdf_page_images(self, page: Any) -> int:
        try:
            page_images = getattr(page, "images", None)
            if page_images is not None:
                return len(page_images)
        except Exception:
            pass

        try:
            resources = page.get("/Resources")
            if resources is None:
                return 0
            resources_obj = resources.get_object() if hasattr(resources, "get_object") else resources
            xobject = resources_obj.get("/XObject") if isinstance(resources_obj, dict) else None
            if xobject is None:
                return 0
            xobject_obj = xobject.get_object() if hasattr(xobject, "get_object") else xobject

            count = 0
            if isinstance(xobject_obj, dict):
                for obj in xobject_obj.values():
                    target = obj.get_object() if hasattr(obj, "get_object") else obj
                    subtype = target.get("/Subtype") if hasattr(target, "get") else None
                    if subtype == "/Image":
                        count += 1
            return count
        except Exception:
            return 0

    def _warm_up_image_components(
        self,
        image_path: Path,
        warmed_components: list[str],
        failed_components: list[str],
    ) -> None:
        strategy = self._image_processing_strategy

        if strategy in {"smoldocling", "hybrid"}:
            self._ensure_image_vlm_converter()
            if self._image_vlm_converter is None:
                failed_components.append("image_smoldocling")
            else:
                output = self._convert_with_docling(self._image_vlm_converter, image_path).strip()
                if output:
                    warmed_components.append("image_smoldocling")
                else:
                    failed_components.append("image_smoldocling")

        if strategy in {"rapid", "hybrid"}:
            self._ensure_image_rapidocr_converter()
            if self._image_rapidocr_converter is None:
                failed_components.append("image_rapidocr")
            else:
                output = self._convert_with_docling(self._image_rapidocr_converter, image_path).strip()
                if output:
                    warmed_components.append("image_rapidocr")
                else:
                    failed_components.append("image_rapidocr")

    def _warm_up_pdf_components(
        self,
        pdf_path: Path,
        warmed_components: list[str],
        failed_components: list[str],
    ) -> None:
        if self._pdf_layout_model_strategy in {"auto", "smoldocling"}:
            self._ensure_pdf_vlm_converter()
            if self._pdf_vlm_converter is None:
                failed_components.append("pdf_smoldocling")
            else:
                output = self._convert_with_docling(self._pdf_vlm_converter, pdf_path).strip()
                if output:
                    warmed_components.append("pdf_smoldocling")
                else:
                    failed_components.append("pdf_smoldocling")

        strategy = self._pdf_ocr_strategy

        if strategy in {"rapid", "hybrid"}:
            self._ensure_pdf_rapidocr_converter()
            if self._pdf_rapidocr_converter is None:
                failed_components.append("pdf_rapidocr")
            else:
                output = self._convert_with_docling(self._pdf_rapidocr_converter, pdf_path).strip()
                if output:
                    warmed_components.append("pdf_rapidocr")
                else:
                    failed_components.append("pdf_rapidocr")

        if strategy in {"easy", "hybrid"}:
            self._ensure_pdf_easyocr_converter()
            if self._pdf_easyocr_converter is None:
                failed_components.append("pdf_easyocr")
            else:
                output = self._convert_with_docling(self._pdf_easyocr_converter, pdf_path).strip()
                if output:
                    warmed_components.append("pdf_easyocr")
                else:
                    failed_components.append("pdf_easyocr")

    def _create_warmup_image(self, image_path: Path) -> bool:
        try:
            from PIL import Image, ImageDraw

            image = Image.new("RGB", (640, 220), "white")
            draw = ImageDraw.Draw(image)
            draw.text((24, 80), "OCR warmup sample 2026", fill="black")
            image.save(image_path)
            return True
        except Exception:
            return False

    def _create_warmup_pdf(self, pdf_path: Path) -> bool:
        try:
            from PIL import Image, ImageDraw

            image = Image.new("RGB", (640, 220), "white")
            draw = ImageDraw.Draw(image)
            draw.text((24, 80), "PDF OCR warmup sample 2026", fill="black")
            image.save(pdf_path, "PDF")
            return True
        except Exception:
            return False

    def _convert_image_with_policy(self, file_path: Path) -> tuple[str, str]:
        strategy = self._image_processing_strategy

        if strategy in {"rapid", "hybrid"}:
            self._ensure_image_rapidocr_converter()
            rapid_text = self._convert_with_docling(self._image_rapidocr_converter, file_path).strip()
            if rapid_text:
                return rapid_text, "image_rapidocr"
            if strategy == "rapid":
                raise RuntimeError("Image OCR failed with RapidOCR")

        if strategy in {"smoldocling", "hybrid"}:
            self._ensure_image_vlm_converter()
            if self._image_vlm_converter is not None:
                vlm_text = self._convert_with_docling(self._image_vlm_converter, file_path).strip()
                if vlm_text:
                    return vlm_text, "image_smoldocling"
            if strategy == "smoldocling":
                raise RuntimeError(
                    "Image conversion failed with SmolDocling VLM. Check VLM dependencies and GPU runtime settings."
                )

        raise RuntimeError("Image conversion failed with RapidOCR and SmolDocling VLM")

    def _fallback_engine_for_suffix(self, suffix: str) -> str:
        mapping = {
            ".txt": "fallback_text",
            ".md": "fallback_text",
            ".csv": "fallback_text",
            ".json": "fallback_text",
            ".pdf": "fallback_pdf",
            ".docx": "fallback_docx",
            ".xlsx": "fallback_xlsx",
            ".xls": "fallback_xls",
            ".ppt": "fallback_ppt",
            ".pptx": "fallback_pptx",
        }
        return mapping.get(suffix, "fallback_unknown")

    def _convert_pdf_with_ocr_policy(self, default_text: str, file_path: Path) -> str:
        if not self._should_try_pdf_ocr(default_text):
            self._pdf_ocr_engine = "pypdf"
            return default_text

        strategy = self._pdf_ocr_strategy
        best_text = default_text
        best_engine = "pypdf" if default_text else "none"

        if strategy == "off":
            self._pdf_ocr_engine = best_engine
            return best_text

        if strategy in {"rapid", "hybrid"}:
            self._ensure_pdf_rapidocr_converter()
            rapid_text = self._convert_with_docling(self._pdf_rapidocr_converter, file_path).strip()
            if self._is_better_pdf_text(best_text, rapid_text):
                best_text = rapid_text
                best_engine = "rapidocr"

            if strategy == "rapid":
                self._pdf_ocr_engine = best_engine
                return best_text

        if strategy in {"easy", "hybrid"}:
            self._ensure_pdf_easyocr_converter()
            easy_text = self._convert_with_docling(self._pdf_easyocr_converter, file_path).strip()
            if self._is_better_pdf_text(best_text, easy_text):
                best_text = easy_text
                best_engine = "easyocr"

        self._pdf_ocr_engine = best_engine
        return best_text

    def _extract_pdf_text_fast(self, file_path: Path) -> str:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            return "\n".join((page.extract_text() or "") for page in reader.pages).strip()
        except Exception:
            return ""

    def _should_try_pdf_ocr(self, default_text: str) -> bool:
        if not default_text:
            return True
        return len(default_text.replace("\n", " ").strip()) < self._pdf_ocr_min_chars

    def _is_better_pdf_text(self, current_text: str, candidate_text: str) -> bool:
        candidate_len = len(candidate_text.replace("\n", " ").strip())
        if candidate_len == 0:
            return False

        current_len = len(current_text.replace("\n", " ").strip())
        if current_len == 0:
            return True

        if current_len < self._pdf_ocr_min_chars:
            return candidate_len >= current_len + 40

        return candidate_len >= int(current_len * 1.25)

    def _ensure_converter(self) -> None:
        if self._converter_initialized:
            return
        self._converter_initialized = True

        try:
            from docling.document_converter import DocumentConverter

            self._converter = DocumentConverter()
        except Exception:
            self._converter = None

    def _ensure_pdf_rapidocr_converter(self) -> None:
        if self._pdf_rapidocr_initialized:
            return
        self._pdf_rapidocr_initialized = True
        self._pdf_rapidocr_converter = self._build_pdf_ocr_converter_with_rapidocr()

    def _ensure_pdf_easyocr_converter(self) -> None:
        if self._pdf_easyocr_initialized:
            return
        self._pdf_easyocr_initialized = True
        self._pdf_easyocr_converter = self._build_pdf_ocr_converter_with_easyocr()

    def _ensure_pdf_vlm_converter(self) -> None:
        if self._pdf_vlm_initialized:
            return
        self._pdf_vlm_initialized = True
        self._pdf_vlm_converter = self._build_pdf_vlm_converter()

    def _ensure_image_vlm_converter(self) -> None:
        if self._image_vlm_initialized:
            return
        self._image_vlm_initialized = True
        self._image_vlm_converter = self._build_image_vlm_converter()

    def _ensure_image_rapidocr_converter(self) -> None:
        if self._image_rapidocr_initialized:
            return
        self._image_rapidocr_initialized = True
        self._image_rapidocr_converter = self._build_image_rapidocr_converter()

    def _build_image_rapidocr_converter(self) -> Any | None:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions, RapidOcrOptions
            from docling.document_converter import DocumentConverter, ImageFormatOption

            pipeline_options = PdfPipelineOptions()
            self._configure_fast_ocr_pipeline(pipeline_options)
            pipeline_options.ocr_options = RapidOcrOptions()

            if hasattr(pipeline_options.ocr_options, "force_full_page_ocr"):
                pipeline_options.ocr_options.force_full_page_ocr = True

            return DocumentConverter(
                format_options={
                    InputFormat.IMAGE: ImageFormatOption(pipeline_options=pipeline_options),
                }
            )
        except Exception:
            return None

    def _build_image_vlm_converter(self) -> Any | None:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import VlmConvertOptions, VlmPipelineOptions
            from docling.document_converter import DocumentConverter, ImageFormatOption
            from docling.models.inference_engines.vlm.transformers_engine import TransformersVlmEngineOptions
            from docling.pipeline.vlm_pipeline import VlmPipeline

            accelerator_device = self._resolve_accelerator_device(self._image_vlm_device)

            engine_options = TransformersVlmEngineOptions(
                device=accelerator_device,
                load_in_8bit=False,
            )
            vlm_options = VlmConvertOptions.from_preset(
                self._image_vlm_preset,
                engine_options=engine_options,
            )
            pipeline_options = VlmPipelineOptions(vlm_options=vlm_options)

            return DocumentConverter(
                format_options={
                    InputFormat.IMAGE: ImageFormatOption(
                        pipeline_cls=VlmPipeline,
                        pipeline_options=pipeline_options,
                    ),
                }
            )
        except Exception:
            return None

    def _build_pdf_vlm_converter(self) -> Any | None:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import VlmConvertOptions, VlmPipelineOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption
            from docling.models.inference_engines.vlm.transformers_engine import TransformersVlmEngineOptions
            from docling.pipeline.vlm_pipeline import VlmPipeline

            accelerator_device = self._resolve_accelerator_device(self._pdf_vlm_device)

            engine_options = TransformersVlmEngineOptions(
                device=accelerator_device,
                load_in_8bit=False,
            )
            vlm_options = VlmConvertOptions.from_preset(
                self._pdf_vlm_preset,
                engine_options=engine_options,
            )
            pipeline_options = VlmPipelineOptions(vlm_options=vlm_options)

            return DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(
                        pipeline_cls=VlmPipeline,
                        pipeline_options=pipeline_options,
                    ),
                }
            )
        except Exception:
            return None

    def _resolve_accelerator_device(self, device_name: str) -> Any:
        from docling.datamodel.accelerator_options import AcceleratorDevice

        device_map = {
            "auto": AcceleratorDevice.AUTO,
            "cpu": AcceleratorDevice.CPU,
            "cuda": AcceleratorDevice.CUDA,
            "mps": AcceleratorDevice.MPS,
            "xpu": AcceleratorDevice.XPU,
        }
        return device_map.get(device_name, AcceleratorDevice.AUTO)

    def _configure_fast_ocr_pipeline(self, pipeline_options: Any) -> None:
        if hasattr(pipeline_options, "do_ocr"):
            pipeline_options.do_ocr = True
        if hasattr(pipeline_options, "do_table_structure"):
            pipeline_options.do_table_structure = False
        if hasattr(pipeline_options, "do_picture_classification"):
            pipeline_options.do_picture_classification = False
        if hasattr(pipeline_options, "do_picture_description"):
            pipeline_options.do_picture_description = False
        if hasattr(pipeline_options, "do_chart_extraction"):
            pipeline_options.do_chart_extraction = False
        if hasattr(pipeline_options, "do_code_enrichment"):
            pipeline_options.do_code_enrichment = False
        if hasattr(pipeline_options, "do_formula_enrichment"):
            pipeline_options.do_formula_enrichment = False

    def _build_pdf_ocr_converter_with_easyocr(self) -> Any | None:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import EasyOcrOptions, PdfPipelineOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption

            pipeline_options = PdfPipelineOptions()
            self._configure_fast_ocr_pipeline(pipeline_options)
            pipeline_options.ocr_options = EasyOcrOptions(
                lang=["ko", "en"],
                use_gpu=self._pdf_vlm_device == "cuda",
            )

            if hasattr(pipeline_options.ocr_options, "force_full_page_ocr"):
                pipeline_options.ocr_options.force_full_page_ocr = True

            return DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
                }
            )
        except Exception:
            return None

    def _build_pdf_ocr_converter_with_rapidocr(self) -> Any | None:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions, RapidOcrOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption

            pipeline_options = PdfPipelineOptions()
            self._configure_fast_ocr_pipeline(pipeline_options)
            pipeline_options.ocr_options = RapidOcrOptions()

            if hasattr(pipeline_options.ocr_options, "force_full_page_ocr"):
                pipeline_options.ocr_options.force_full_page_ocr = True

            return DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
                }
            )
        except Exception:
            return None

    def _convert_with_docling(self, converter: Any, file_path: Path) -> str:
        if converter is None:
            return ""

        try:
            result = converter.convert(str(file_path))
            if hasattr(result, "document") and hasattr(result.document, "export_to_markdown"):
                markdown = result.document.export_to_markdown()
                return str(markdown or "")
            return str(result or "")
        except Exception:
            return ""

    def _fallback_convert(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()

        if suffix in {".txt", ".md", ".csv", ".json"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")

        if suffix == ".pdf":
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(file_path))
                extracted = "\n".join((page.extract_text() or "") for page in reader.pages).strip()
                if extracted:
                    return extracted
                raise RuntimeError(
                    "Converted text is empty. This PDF appears to be image/scanned without a selectable text layer. OCR is required."
                )
            except RuntimeError:
                raise
            except Exception as exc:
                raise RuntimeError("PDF conversion failed") from exc

        if suffix == ".docx":
            try:
                from docx import Document

                document = Document(str(file_path))
                return "\n".join(paragraph.text for paragraph in document.paragraphs)
            except Exception as exc:
                raise RuntimeError("DOCX conversion failed") from exc

        if suffix == ".xlsx":
            try:
                from openpyxl import load_workbook

                workbook = load_workbook(str(file_path), read_only=True, data_only=True)
                lines: list[str] = []

                for sheet in workbook.worksheets:
                    lines.append(f"[Sheet] {sheet.title}")
                    for row in sheet.iter_rows(values_only=True):
                        values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                        if values:
                            lines.append(" | ".join(values))

                return "\n".join(lines)
            except Exception as exc:
                raise RuntimeError("XLSX conversion failed") from exc

        if suffix == ".xls":
            try:
                import xlrd

                workbook = xlrd.open_workbook(str(file_path))
                lines: list[str] = []

                for sheet in workbook.sheets():
                    lines.append(f"[Sheet] {sheet.name}")
                    for row_idx in range(sheet.nrows):
                        row_values = [str(cell).strip() for cell in sheet.row_values(row_idx) if str(cell).strip()]
                        if row_values:
                            lines.append(" | ".join(row_values))

                return "\n".join(lines)
            except Exception as exc:
                raise RuntimeError("XLS conversion failed") from exc

        if suffix in {".ppt", ".pptx"}:
            try:
                from pptx import Presentation

                presentation = Presentation(str(file_path))
                lines: list[str] = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            lines.append(shape.text)
                return "\n".join(lines)
            except Exception as exc:
                raise RuntimeError("PPT conversion failed") from exc

        raise RuntimeError("No available converter for this file type")
