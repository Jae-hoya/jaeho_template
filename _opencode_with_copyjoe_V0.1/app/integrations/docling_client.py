from pathlib import Path


class DoclingClient:
    def __init__(self) -> None:
        self._converter = None

        try:
            from docling.document_converter import DocumentConverter

            self._converter = DocumentConverter()
        except Exception:
            self._converter = None

    def convert_to_text(self, file_path: Path) -> str:
        if self._converter is not None:
            try:
                result = self._converter.convert(str(file_path))
                if hasattr(result, "document") and hasattr(result.document, "export_to_markdown"):
                    return result.document.export_to_markdown()
                return str(result)
            except Exception:
                pass

        return self._fallback_convert(file_path)

    def _fallback_convert(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()

        if suffix in {".txt", ".md", ".csv", ".json"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")

        if suffix == ".pdf":
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(file_path))
                return "\n".join((page.extract_text() or "") for page in reader.pages)
            except Exception as exc:
                raise RuntimeError("PDF conversion failed") from exc

        if suffix == ".docx":
            try:
                from docx import Document

                document = Document(str(file_path))
                return "\n".join(paragraph.text for paragraph in document.paragraphs)
            except Exception as exc:
                raise RuntimeError("DOCX conversion failed") from exc

        if suffix == ".xlsx":
            try:
                from openpyxl import load_workbook

                workbook = load_workbook(str(file_path), read_only=True, data_only=True)
                lines: list[str] = []

                for sheet in workbook.worksheets:
                    lines.append(f"[Sheet] {sheet.title}")
                    for row in sheet.iter_rows(values_only=True):
                        values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                        if values:
                            lines.append(" | ".join(values))

                return "\n".join(lines)
            except Exception as exc:
                raise RuntimeError("XLSX conversion failed") from exc

        if suffix == ".xls":
            try:
                import xlrd

                workbook = xlrd.open_workbook(str(file_path))
                lines: list[str] = []

                for sheet in workbook.sheets():
                    lines.append(f"[Sheet] {sheet.name}")
                    for row_idx in range(sheet.nrows):
                        row_values = [str(cell).strip() for cell in sheet.row_values(row_idx) if str(cell).strip()]
                        if row_values:
                            lines.append(" | ".join(row_values))

                return "\n".join(lines)
            except Exception as exc:
                raise RuntimeError("XLS conversion failed") from exc

        if suffix in {".ppt", ".pptx"}:
            try:
                from pptx import Presentation

                presentation = Presentation(str(file_path))
                lines: list[str] = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            lines.append(shape.text)
                return "\n".join(lines)
            except Exception as exc:
                raise RuntimeError("PPT conversion failed") from exc

        raise RuntimeError("No available converter for this file type")
